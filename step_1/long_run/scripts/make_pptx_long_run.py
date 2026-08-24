"""
Long-run LinkedIn PDF brief — 3 slides, A4 Landscape.

  python step_1/long_run/scripts/make_pptx_long_run.py

Charts must exist first:  python step_1/long_run/scripts/viz_long_run.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt, Emu

STEP_ROOT = Path(__file__).resolve().parents[2]   # step_1/
LONG_RUN = STEP_ROOT / "long_run"
FIGURES = LONG_RUN / "exports" / "figures"
OUT = LONG_RUN / "exports"

SLIDE_W, SLIDE_H = Cm(29.7), Cm(21.0)

INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WARM = RGBColor(0xD9, 0x77, 0x06)
BAR_C = RGBColor(0x25, 0x63, 0xEB)
LINE_C = RGBColor(0xD1, 0xD5, 0xDB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG = RGBColor(0xF3, 0xF4, 0xF6)

FONT = "Calibri"
ML = Cm(1.2)
CW = SLIDE_W - 2 * ML
COL_W = (CW - Cm(1.0)) // 2

FOOTER = (
    "Sources: OCSTAT (notarial transactions, Canton of Geneva);"
    " W\u00fcest Partner (hedonic index, not redistributed);"
    " BIS/FRED; BFS IMPI"
    " \u00b7 Author\u2019s calculations"
)


def _bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(0.25))
    s.fill.solid()
    s.fill.fore_color.rgb = BAR_C
    s.line.fill.background()


def _txt(slide, left, top, w, h, text, *, sz=12, bold=False, color=INK,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = align
    return tb


def _bullets(slide, left, top, w, h, lines, *, sz=9.5, color=MUTED,
             spacing=Pt(5)):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
    return tb


def _img(slide, name, left, top, w):
    path = FIGURES / name
    if not path.exists():
        print(f"  WARNING: missing {name}")
        return None
    return slide.shapes.add_picture(str(path), left, top, width=w)


def _hline(slide, left, top, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Cm(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = LINE_C
    s.line.fill.background()


def _footer(slide):
    _txt(slide, ML, Cm(19.8), CW, Cm(0.5), FOOTER, sz=7, color=MUTED)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _cell_fmt(cell, text, *, sz=8, bold=False, color=INK, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)


def _add_growth_table(slide, left, top, w):
    """Comparative growth table: apartments + houses + inflation."""
    rows_data = [
        ("APARTMENTS (PPE, CHF/m\u00b2)", "", "", "", ""),
        ("Geneva city", "OCSTAT", "264", "+2.9%", "+1.8%"),
        ("Geneva canton", "OCSTAT", "234", "+2.5%", "+1.4%"),
        ("Lake Geneva region", "W\u00fcest", "278", "+3.0%", "+1.9%"),
        ("Switzerland", "W\u00fcest", "235", "+2.5%", "+1.4%"),
        ("HOUSES (CHF/object)", "", "", "", ""),
        ("Geneva canton", "OCSTAT", "232", "+2.4%", "+1.4%"),
        ("Lake Geneva region", "W\u00fcest", "265", "+2.8%", "+1.8%"),
        ("Switzerland", "W\u00fcest", "231", "+2.5%", "+1.4%"),
        ("INFLATION & COSTS", "", "", "", ""),
        ("Geneva CPI", "OCSTAT", "144", "+1.0%", "\u2014"),
        ("Switzerland CPI", "BFS", "140", "+1.0%", "\u2014"),
        ("Geneva construction cost", "OCSTAT", "160", "+1.3%", "\u2014"),
    ]
    headers = ["Geography", "Source", "Index\n(1990=100)", "Nominal\n/year", "Real\n/year"]
    n_rows = len(rows_data) + 1
    n_cols = 5

    col_widths = [Cm(5.2), Cm(2.4), Cm(2.2), Cm(2.2), Cm(2.0)]
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, w, Cm(7.5))
    tbl = table_shape.table

    for ci in range(n_cols):
        tbl.columns[ci].width = col_widths[ci]

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        _cell_fmt(cell, h, sz=8, bold=True, color=WHITE)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT

    section_rows = {0, 5, 9}
    for ri, row in enumerate(rows_data):
        is_section = ri in section_rows
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if is_section and ci == 0:
                _cell_fmt(cell, val, sz=7.5, bold=True, color=ACCENT)
            elif is_section:
                _cell_fmt(cell, val, sz=7.5, color=MUTED)
            else:
                a = PP_ALIGN.RIGHT if ci >= 2 else PP_ALIGN.LEFT
                _cell_fmt(cell, val, sz=8, color=INK, align=a)
            if is_section:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xEF, 0xF6, 0xFF)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    WHITE if (ri % 2 == 1) else HEADER_BG
                )

    tbl.first_row = False
    tbl.horz_banding = False
    return table_shape


# ── Slide 1: Long-term trends ────────────────────────────────────────

def page_trends(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "Geneva \u00b7 Long-Term Price Dynamics \u00b7 1990\u20132025",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "Data Brief #2",
         sz=14, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "Three independent sources, one conclusion:"
         " Geneva real estate roughly tripled since 1990.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(2.3)

    _img(sl, "lr_01_normalized_trend.png", lx, y_chart, COL_W)
    _img(sl, "lr_06_houses_trend.png", rx, y_chart, COL_W)

    y_ins = Cm(12.8)
    _bullets(sl, ML, y_ins, CW, Cm(6.5), [
        "\u25b8 Apartments: city +2.9%/yr nominal (+1.8% real),"
        " canton +2.5%/yr (+1.4% real) over 34 years."
        " Houses: canton +2.4%/yr, Lake Geneva region +2.8%/yr."
        " Prices roughly double every 25\u201328 years in nominal terms.",

        "\u25b8 Official OCSTAT transactions and W\u00fcest hedonic index"
        " show the same trajectory"
        " (correlation 0.99 on levels, 74\u201390% direction agreement)."
        " The trend is not an artifact of one measurement method.",

        "\u25b8 Geneva follows the Swiss cycle but amplifies it:"
        " deeper corrections, steeper rises."
        " CPI grew just +1.0%/yr \u2014 real estate outpaced inflation"
        " by +1.4\u20131.9 percentage points annually.",
    ], sz=9.5)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


# ── Slide 2: Drawdowns ───────────────────────────────────────────────

def page_drawdowns(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "Historical Drawdowns \u00b7 What Happens After a Peak?",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "Apartments (PPE)",
         sz=10, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "Two completed correction cycles \u2014"
         " both fully recovered, each new peak higher than the last.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    _img(sl, "lr_05_drawdowns.png", ML, Cm(2.3), CW)

    y_ins = Cm(13.8)
    _bullets(sl, ML, y_ins, CW, Cm(6.0), [
        "\u25b8 Worst case: Geneva city fell \u221225% from peak"
        " and took 14 years to recover (1990\u20132004)."
        " Canton: \u221221%, 16 years."
        " Switzerland nationally: only \u221210%"
        " \u2014 Geneva amplifies the cycle 2\u00d7.",

        "\u25b8 The second cycle (2012\u20132021) was milder:"
        " \u221210 to \u221212%, recovery in 6\u20139 years."
        " Shorter, shallower \u2014 potentially a more"
        " resilient market or different rate environment.",

        "\u25b8 After every drawdown, the new peak substantially"
        " exceeded the previous one."
        " Long-term holders were always rewarded"
        " \u2014 but the holding period matters.",
    ], sz=9.5)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


# ── Slide 3: Inflation, costs & summary table ────────────────────────

def page_costs_summary(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "Inflation, Construction Costs & Growth Summary",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "All indices: 1990 = 100",
         sz=10, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "Real estate outpaced inflation by +1.4\u20131.9 p.p./year."
         " Construction costs surged after 2020.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)

    _img(sl, "lr_07_prices_vs_costs.png", lx, Cm(2.3), COL_W)

    _add_growth_table(sl, rx, Cm(2.4), COL_W)

    y_ins = Cm(13.2)
    _bullets(sl, ML, y_ins, CW, Cm(6.0), [
        "\u25b8 Geneva CPI +1.0%/yr, construction costs +1.3%/yr"
        " (with a sharp acceleration after 2020)."
        " Real estate nominal growth of +2.4\u20133.0%/yr"
        " translates to +1.4\u20131.9%/yr above inflation.",

        "\u25b8 Construction cost surge (+18% in 4 years)"
        " supports the new-build price premium"
        " and limits future supply growth"
        " \u2014 a structural tailwind for existing property values.",

        "\u25b8 Data: OCSTAT (official notarial transactions),"
        " W\u00fcest Partner (hedonic quality-adjusted index),"
        " BIS/FRED, BFS IMPI."
        " Full methodology and code: GitHub.",
    ], sz=9.5)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


# ── Slide 4: Implied apartment size ───────────────────────────────────

def page_implied_size(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "What Size Apartment Is Actually Selling?",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "Geneva canton \u00b7 2006\u20132024",
         sz=10, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "Proxy: median object price \u00f7 median CHF/m\u00b2"
         " \u2014 the implied typical floor area of a sold apartment.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    _img(sl, "lr_08_implied_size.png", ML, Cm(2.3), CW)

    y_ins = Cm(13.2)
    _bullets(sl, ML, y_ins, CW, Cm(6.0), [
        "\u25b8 Existing apartments: ~120 m\u00b2, remarkably stable"
        " (+0.1%/yr over 18 years)."
        " The composition of what sells has barely changed"
        " \u2014 a consistent market for 3\u20134 room flats.",

        "\u25b8 New apartments: ~140 m\u00b2 in 2024, trending upward"
        " since 2014 (from ~125 m\u00b2)."
        " Developers are building larger, higher-value units."
        " Early years (2006\u20132012) show high volatility"
        " due to small sample sizes.",

        "\u25b8 New apartments are ~17% larger than existing."
        " This is a proxy (ratio of two medians,"
        " not median floor area directly),"
        " but the direction is consistent with market observations.",
    ], sz=9.5)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page_trends(prs)
    page_drawdowns(prs)
    page_costs_summary(prs)
    page_implied_size(prs)

    out = OUT / "Geneva_Long_Run_Brief.pptx"
    prs.save(str(out))
    print(f"Saved -> {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
