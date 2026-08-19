"""
Build executive brief as PowerPoint A4 Landscape → export as PDF.

  python scripts/make_pptx.py

Charts must exist first:  python scripts/viz.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

STEP_ROOT = Path(__file__).resolve().parents[1]   # step_1/
FIGURES = STEP_ROOT / "exports" / "figures"
OUT = STEP_ROOT / "exports"

SLIDE_W, SLIDE_H = Cm(29.7), Cm(21.0)

INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WARM = RGBColor(0xD9, 0x77, 0x06)
KPI_BG = RGBColor(0xF3, 0xF4, 0xF6)
BAR_C = RGBColor(0x25, 0x63, 0xEB)
LINE_C = RGBColor(0xD1, 0xD5, 0xDB)

FONT = "Calibri"
ML = Cm(1.2)
CW = SLIDE_W - 2 * ML
COL_W = (CW - Cm(1.0)) // 2


def _bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(0.25))
    s.fill.solid()
    s.fill.fore_color.rgb = BAR_C
    s.line.fill.background()


def _txt(slide, left, top, w, h, text, *, sz=12, bold=False, color=INK,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = align
    return tb


def _bullets(slide, left, top, w, h, lines, *, sz=9, color=MUTED,
             spacing=Pt(3)):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
    return tb


def _img(slide, name, left, top, w):
    path = FIGURES / name
    if not path.exists():
        print(f"  WARNING: missing {name}")
        return None
    return slide.shapes.add_picture(str(path), left, top, width=w)


def _hline(slide, left, top, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Cm(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = LINE_C
    s.line.fill.background()


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Slide 1: Snapshot Q2 2026 (two-column) ──────────────────────────────

def slide_snapshot(prs):
    sl = _blank(prs)
    _bar(sl)

    # ── Header ──
    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Geneva Canton \u00b7 Residential Market",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "Official OCSTAT transactions \u00b7 Q2 2026",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    # ── Left column: Chart 09 + insights ──
    lx = ML
    y_chart = Cm(1.5)
    _img(sl, "09_q_latest_avg_ticket.png", lx, y_chart, COL_W)

    y_ins = Cm(11.8)
    _txt(sl, lx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, lx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 A house costs on average 2\u00d7 more"
        " than an apartment (3.6M vs 1.8M CHF).",

        "\u25b8 New apartment average rose +56% YoY \u2014"
        " driven by several large projects closing,"
        " not a broad market-wide price increase.",

        "\u25b8 Existing house average (3.6M) is pushed up"
        " by luxury transactions in premium communes"
        " (Cologny, Vandoeuvres).",

        "\u25b8 New and existing apartment prices"
        " are similar (~1.8\u20132.0M per transaction).",
    ], sz=8.5)

    # ── Right column: Chart 14 + insights ──
    rx = ML + COL_W + Cm(1.0)
    _img(sl, "14_q_latest_existing_range.png", rx, y_chart, COL_W)

    _txt(sl, rx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, rx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Half of house buyers pay between 1.6M and"
        " 3.2M CHF \u2014 prices vary significantly"
        " by commune and property type.",

        "\u25b8 Apartments are more predictable:"
        " 75% of transactions fall within 1.0\u20131.9M.",

        "\u25b8 Apartment P75 (~1.9M) meets house P25 (~1.6M)"
        " \u2014 at this price, a buyer chooses between"
        " a large apartment or a modest house.",

        "\u25b8 House median stable at 2.10M (\u22122.1% YoY),"
        " apartment median 1.40M (+3.5% YoY).",
    ], sz=8.5)

    # ── Footer ──
    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05 \u00b7 Canton of Geneva \u00b7 YoY vs Q2 2025"
         " \u00b7 Chart 1: average ticket (value/n), not median"
         " \u00b7 Chart 2: official OCSTAT medians + P25\u2013P75",
         sz=7, color=MUTED)


# ── Slide 2: CHF/m² canton overview ───────────────────────────────────

def slide_m2_canton(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Price per m\u00b2 \u00b7 Canton Overview",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "OCSTAT \u00b7 Notarial transactions \u00b7 Annual 2024",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(1.5)
    y_ins = Cm(11.8)

    # ── Left: 06b canton new vs existing ──
    _img(sl, "06b_apartments_2024_quartiles.png", lx, y_chart, COL_W)

    _txt(sl, lx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, lx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Existing apartments (10,853 CHF/m\u00b2) are 5.5%"
        " more expensive than new (10,284) at canton level"
        " \u2014 counterintuitive.",

        "\u25b8 Secondary market is 2.1\u00d7 larger:"
        " 828 existing vs 395 new transactions.",

        "\u25b8 Both grew ~2.4% YoY \u2014 market moves"
        " uniformly, not by segment.",

        "\u25b8 Existing P25\u2013P75 range (8.9\u201312.9K)"
        " is 1.7\u00d7 wider than new (9.3\u201311.6K)"
        " \u2014 more diverse locations.",
    ], sz=8.5)

    # ── Right: 04c dumbbell ──
    _img(sl, "04c_new_vs_existing_commune.png", rx, y_chart, COL_W)

    _txt(sl, rx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, rx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Within each commune, new costs 17\u201329%"
        " MORE than existing \u2014 the opposite"
        " of the canton-level picture.",

        "\u25b8 This is a composition effect:"
        " new builds concentrate in affordable communes"
        " (Thônex, Meyrin, Bernex),"
        " while existing stock trades everywhere"
        " including premium locations.",

        "\u25b8 Bernex has the highest new-build premium"
        " (+29%), likely single large project at n=15.",

        "\u25b8 Practical: \u201cnew is cheaper\u201d is a geographic"
        " illusion \u2014 same neighborhood, same m\u00b2,"
        " new always costs more.",
    ], sz=8.5)

    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05.1.4.03 \u00b7 Canton of Geneva"
         " \u00b7 All prices = completed notarial transactions,"
         " not listings or appraisals"
         " \u00b7 CHF per square meter, free-market apartments (PPE)",
         sz=7, color=MUTED)


# ── Slide 3: CHF/m² by commune ────────────────────────────────────────

def slide_m2_communes(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Price per m\u00b2 \u00b7 Communes",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "OCSTAT \u00b7 Notarial transactions \u00b7 Annual 2024",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(1.5)
    y_ins = Cm(11.8)

    # ── Left: 04 existing by commune ──
    _img(sl, "04_ppe_range.png", lx, y_chart, COL_W)

    _txt(sl, lx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, lx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 1.8\u00d7 price gap: Vernier (7,182 CHF/m\u00b2)"
        " vs Genève (12,597). Same canton,"
        " 75% more living space for same budget.",

        "\u25b8 Only 4 communes above canton median"
        " (10,853): Genève, Carouge,"
        " Cologny, Chêne-Bougeries.",

        "\u25b8 Carouge: strongest growth (+15.9% YoY)"
        " among premium communes.",

        "\u25b8 Veyrier (\u221213.2%) and Plan-les-Ouates"
        " (\u22125.8%) declined \u2014 but at n=18"
        " and n=12, this may be statistical noise.",
    ], sz=8.5)

    # ── Right: 04b new by commune ──
    _img(sl, "04b_ppe_range_new.png", rx, y_chart, COL_W)

    _txt(sl, rx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, rx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Only 5 of 45 communes have enough"
        " new-apartment data (n\u226510)"
        " \u2014 market is extremely concentrated.",

        "\u25b8 Two tiers: affordable (Meyrin, Thônex,"
        " Bernex at 10\u201311K) vs premium"
        " (Genève, Cologny at 15\u201316K).",

        "\u25b8 Genève new +23.1% YoY \u2014 not market"
        " growth but a few luxury project closings"
        " (n=27).",

        "\u25b8 Communes appear and disappear from"
        " the statistics each year \u2014 new-build"
        " market is structurally volatile.",
    ], sz=8.5)

    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05.1.4.03 \u00b7 Canton of Geneva"
         " \u00b7 Unpublished communes (n<10) omitted,"
         " not zero \u00b7 YoY vs 2023",
         sz=7, color=MUTED)


# ── Slide 4: Long-term price trends ───────────────────────────────────

def slide_trends(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Long-term Trends \u00b7 How We Got Here",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "OCSTAT \u00b7 Canton of Geneva \u00b7 Annual",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(1.5)
    y_ins = Cm(11.8)

    # ── Left: 01 PPE trend ──
    _img(sl, "01_ppe_trend.png", lx, y_chart, COL_W)

    _txt(sl, lx, y_ins, COL_W, Cm(0.5),
         "Apartments (CHF/m\u00b2)", sz=9, bold=True, color=ACCENT)
    _bullets(sl, lx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Median roughly doubled since 2006:"
        " from ~5,300 to ~10,900 CHF/m\u00b2.",

        "\u25b8 Fast growth 2006\u20132012 (+80%),"
        " then a 8-year plateau (2012\u20132020),"
        " followed by renewed acceleration.",

        "\u25b8 During the 2012\u20132016 correction, P75"
        " fell 14% while P25 fell only 5%"
        " \u2014 expensive segment corrects faster.",

        "\u25b8 Even during plateaus, prices never fell"
        " below the prior cycle peak"
        " \u2014 no sustained downward correction.",
    ], sz=8.5)

    # ── Right: 20 houses trend ──
    _img(sl, "20_houses_trend.png", rx, y_chart, COL_W)

    _txt(sl, rx, y_ins, COL_W, Cm(0.5),
         "Houses (CHF per dwelling)", sz=9, bold=True, color=WARM)
    _bullets(sl, rx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Median house tripled since the mid-1990s:"
        " from ~0.75M to ~2.2M CHF.",

        "\u25b8 1990\u20132004: decade of stagnation"
        " (flat at ~0.8\u20131.0M)."
        " 2004\u20132012: sharp rally to ~1.8M.",

        "\u25b8 During 2012\u20132017 correction, P75"
        " dropped 21% (2.65M \u2192 2.10M) while P25"
        " fell only 7% \u2014 luxury corrects first.",

        "\u25b8 Post-2020: P75 surged to 3.3M (+58%)"
        " then stabilized at 3.1M, while median"
        " grew only 22% \u2014 premium segment"
        " is more volatile in both directions.",
    ], sz=8.5)

    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05 \u00b7 Canton of Geneva"
         " \u00b7 Left: existing PPE, free market, CHF/m\u00b2"
         " \u00b7 Right: detached houses, CHF per dwelling"
         " \u00b7 Band = P25\u2013P75 of transactions",
         sz=7, color=MUTED)


# ── Slide 5: Market resilience + CAGR ─────────────────────────────────

def slide_resilience(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Market Resilience \u00b7 Shocks & Growth by Commune",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "OCSTAT \u00b7 Canton of Geneva \u00b7 Annual",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(1.5)
    y_ins = Cm(11.8)

    # ── Left: 17 market shocks ──
    _img(sl, "17_market_shocks.png", lx, y_chart, COL_W)

    _txt(sl, lx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, lx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 2008 financial crisis: apartment prices"
        " peaked at +15% YoY growth,"
        " then corrected to \u22123% \u2014"
        " a mild and short-lived dip.",

        "\u25b8 COVID 2020: transaction volumes dropped"
        " ~20%, but prices barely moved"
        " \u2014 demand absorbed the shock.",

        "\u25b8 2022 interest-rate hikes: volumes fell"
        " but prices continued to rise."
        " Geneva prices are \u201csticky downward.\u201d",

        "\u25b8 Key pattern: volumes absorb shocks,"
        " prices remain resilient."
        " No sustained decline in 18 years.",
    ], sz=8.5)

    # ── Right: 21 quarterly corridor ──
    _img(sl, "21_quarterly_corridor.png", rx, y_chart, COL_W)

    _txt(sl, rx, y_ins, COL_W, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, rx, y_ins + Cm(0.5), COL_W, Cm(7.0), [
        "\u25b8 Quarterly median fluctuates \u00b110\u201315%"
        " due to composition (which properties"
        " sold), not actual price changes.",

        "\u25b8 Max drawdown: apartments \u221236%"
        " (Q4\u20192008 \u2192 Q1\u20192009), houses \u221216%"
        " (2015\u20132017). Both recovered fully.",

        "\u25b8 P25\u2013P75 corridor width remains stable"
        " across cycles \u2014 the market shifts"
        " as a whole, not just the extremes.",

        "\u25b8 No quarter-over-quarter decline lasted"
        " more than 2\u20133 quarters \u2014 structural"
        " demand keeps prices resilient.",
    ], sz=8.5)

    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05 \u00b7 Canton of Geneva"
         " \u00b7 Left: YoY % changes, annual"
         " \u00b7 Right: quarterly median + P25\u2013P75,"
         " existing dwellings, CHF per object",
         sz=7, color=MUTED)


# ── Slide 6: Map — Apartments CHF/m² ──────────────────────────────────

def slide_map_apartments(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Geography \u00b7 Apartment Prices per m\u00b2",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "OCSTAT \u00b7 Existing PPE \u00b7 Free market \u00b7 2024",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    map_w = Cm(16.0)
    _img(sl, "03b_map_apartments.png", ML, Cm(1.5), map_w)

    rx = ML + map_w + Cm(0.5)
    rw = SLIDE_W - rx - ML

    _txt(sl, rx, Cm(2.0), rw, Cm(0.5),
         "Key observations", sz=9, bold=True, color=ACCENT)
    _bullets(sl, rx, Cm(2.5), rw, Cm(15.0), [
        "\u25b8 The price gradient runs east-west:"
        " premium communes (Gen\u00e8ve, Cologny,"
        " Ch\u00eane-Bougeries) cluster along the lake"
        " shore at 12\u201312.5K CHF/m\u00b2.",

        "\u25b8 Western suburbs (Vernier, Bernex,"
        " Meyrin) offer 30\u201340% lower prices"
        " per m\u00b2 \u2014 same canton, very"
        " different value proposition.",

        "\u25b8 Carouge stands out: +16% YoY with"
        " 12.4K/m\u00b2 \u2014 catching up to the"
        " premium lakeside communes.",

        "\u25b8 Veyrier (\u221213%) and Vernier (\u221211%)"
        " declined YoY, but small sample sizes"
        " (n=18 and n=73) may explain the drop.",

        "\u25b8 Grey communes have fewer than 10"
        " transactions \u2014 OCSTAT does not"
        " publish a median for them.",

        "\u25b8 1.8\u00d7 price range across the canton:"
        " same 60m\u00b2 apartment costs 430K"
        " in Vernier vs 750K in Cologny.",
    ], sz=8.5)

    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05.1.4.03 \u00b7 Canton of Geneva"
         " \u00b7 Notarial transactions, not listings"
         " \u00b7 Grey = n<10 (unpublished)"
         " \u00b7 YoY vs 2023",
         sz=7, color=MUTED)


# ── Slide 7: Map — Houses CHF/dwelling ────────────────────────────────

def slide_map_houses(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(20), Cm(0.7),
         "Geography \u00b7 House Prices per Dwelling",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(10), Cm(0.5), Cm(10), Cm(0.7),
         "OCSTAT \u00b7 Detached houses \u00b7 2024",
         sz=9, color=MUTED, align=PP_ALIGN.RIGHT)

    _hline(sl, ML, Cm(1.2), CW)

    map_w = Cm(16.0)
    _img(sl, "03c_map_houses.png", ML, Cm(1.5), map_w)

    rx = ML + map_w + Cm(0.5)
    rw = SLIDE_W - rx - ML

    _txt(sl, rx, Cm(2.0), rw, Cm(0.5),
         "Key observations", sz=9, bold=True, color=WARM)
    _bullets(sl, rx, Cm(2.5), rw, Cm(15.0), [
        "\u25b8 Cologny dominates at 6.2M CHF/house"
        " \u2014 3.7\u00d7 more than Vernier (1.7M)."
        " A single canton, two different markets.",

        "\u25b8 The premium belt along the lake"
        " (Cologny, Ch.-Bougeries, Vandoeuvres,"
        " Collonge-Bellerive) = 3.5\u20136.2M.",

        "\u25b8 Gen\u00e8ve city: \u221241% YoY (1.9M)"
        " \u2014 misleading, as only n=18"
        " transactions in the city proper."
        " House market is suburban.",

        "\u25b8 Bernex: +23% YoY at 2.0M \u2014"
        " the strongest grower among"
        " affordable suburbs.",

        "\u25b8 More communes publish house data"
        " (15) than apartment CHF/m\u00b2 (13)"
        " \u2014 houses trade in more communes.",

        "\u25b8 Note: prices are per dwelling,"
        " not per m\u00b2. Larger houses cost more"
        " regardless of the location premium.",
    ], sz=8.5)

    _hline(sl, ML, Cm(19.2), CW)
    _txt(sl, ML, Cm(19.4), CW, Cm(0.5),
         "Source: OCSTAT T 05.05.1.3.03 \u00b7 Canton of Geneva"
         " \u00b7 Notarial transactions"
         " \u00b7 Median CHF per dwelling (not per m\u00b2)"
         " \u00b7 YoY vs 2023",
         sz=7, color=MUTED)


def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_snapshot(prs)
    slide_m2_canton(prs)
    slide_m2_communes(prs)
    slide_trends(prs)
    slide_resilience(prs)
    slide_map_apartments(prs)
    slide_map_houses(prs)

    out = OUT / "Geneva_Market_Brief_v8.pptx"
    prs.save(str(out))
    print(f"Saved -> {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
