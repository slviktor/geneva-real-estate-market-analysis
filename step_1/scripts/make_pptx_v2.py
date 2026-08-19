"""
LinkedIn PDF brief — 3 slides, A4 Landscape.

  python scripts/make_pptx_v2.py

Charts must exist first:  python scripts/viz.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

STEP_ROOT = Path(__file__).resolve().parents[1]   # step_1/
FIGURES = STEP_ROOT / "exports" / "figures"
OUT = STEP_ROOT / "exports"

SLIDE_W, SLIDE_H = Cm(29.7), Cm(21.0)

INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
WARM = RGBColor(0xD9, 0x77, 0x06)
BAR_C = RGBColor(0x25, 0x63, 0xEB)
LINE_C = RGBColor(0xD1, 0xD5, 0xDB)

FONT = "Calibri"
ML = Cm(1.2)
CW = SLIDE_W - 2 * ML
COL_W = (CW - Cm(1.0)) // 2

FOOTER = (
    "Source: OCSTAT T 05.05 \u00b7 Official notarial transactions,"
    " Canton of Geneva \u00b7 Author\u2019s calculations"
)


def _bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Cm(0.25))
    s.fill.solid()
    s.fill.fore_color.rgb = BAR_C
    s.line.fill.background()


def _txt(slide, left, top, w, h, text, *, sz=12, bold=False, color=INK,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.font.italic = italic
    p.alignment = align
    return tb


def _bullets(slide, left, top, w, h, lines, *, sz=9, color=MUTED,
             spacing=Pt(4)):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = spacing
    return tb


def _img(slide, name, left, top, w):
    path = FIGURES / name
    if not path.exists():
        print(f"  WARNING: missing {name}")
        return None
    return slide.shapes.add_picture(str(path), left, top, width=w)


def _hline(slide, left, top, w):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Cm(0.02))
    s.fill.solid()
    s.fill.fore_color.rgb = LINE_C
    s.line.fill.background()


def _footer(slide):
    _txt(slide, ML, Cm(19.8), CW, Cm(0.5), FOOTER, sz=7, color=MUTED)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Page 1: Price Level ───────────────────────────────────────────────

def page_price(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "Geneva Canton \u00b7 Residential Market \u00b7 Q2 2026",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "Data Brief",
         sz=14, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "Apartments and houses are two different markets"
         " \u2014 different metrics, different dynamics.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(2.3)

    _img(sl, "09_q_latest_avg_ticket.png", lx, y_chart, COL_W)
    _img(sl, "14_q_latest_existing_range.png", rx, y_chart, COL_W)

    y_ins = Cm(13.0)
    _bullets(sl, ML, y_ins, CW, Cm(6.0), [
        "\u25b8 A house costs on average 2\u00d7 more than an apartment"
        " (3.6M vs 1.8M CHF), but house prices are plateauing"
        " while apartments keep rising.",

        "\u25b8 Apartment P75 (~1.9M) meets house P25 (~1.6M)"
        " \u2014 the decision boundary between a large apartment"
        " and a modest house.",

        "\u25b8 All figures are completed notarial transactions"
        " (OCSTAT), not listings or appraisals."
        " Left chart: average ticket. Right chart: official"
        " OCSTAT medians with P25\u2013P75 range.",
    ], sz=9)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


# ── Page 2: Geography ─────────────────────────────────────────────────

def page_geography(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "Geography \u00b7 Where Does It Matter?",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "Annual 2024",
         sz=10, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "1.8\u00d7 price gap across communes"
         " \u2014 same canton, very different value.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    lx = ML
    rx = ML + COL_W + Cm(1.0)
    y_chart = Cm(2.3)

    _img(sl, "03b_map_apartments.png", lx, y_chart, COL_W)
    _img(sl, "03c_map_houses.png", rx, y_chart, COL_W)

    y_ins = Cm(13.0)
    _bullets(sl, ML, y_ins, CW, Cm(6.0), [
        "\u25b8 Apartments: premium lakeside"
        " (Cologny, Ch\u00eane-Bougeries: 12\u201312.5K CHF/m\u00b2)"
        " vs western suburbs"
        " (Vernier, Bernex: 7\u20139K)."
        " A 60m\u00b2 flat costs 430K in Vernier"
        " vs 750K in Cologny.",

        "\u25b8 Houses: Cologny median (6.2M) is"
        " 3.7\u00d7 Vernier (1.7M)."
        " Most extreme divergence in the canton.",

        "\u25b8 \u201cNew is cheaper\u201d at canton level"
        " is a composition effect \u2014"
        " new builds concentrate in affordable communes."
        " Within each commune, new costs 17\u201329% more.",
    ], sz=9)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


# ── Page 3: Data & Method ─────────────────────────────────────────────

def page_method(prs):
    sl = _blank(prs)
    _bar(sl)

    _txt(sl, ML, Cm(0.5), Cm(22), Cm(0.7),
         "Data & Methodology",
         sz=14, bold=True)
    _txt(sl, SLIDE_W - ML - Cm(8), Cm(0.5), Cm(8), Cm(0.7),
         "Transparency note",
         sz=10, color=MUTED, align=PP_ALIGN.RIGHT)

    _txt(sl, ML, Cm(1.3), CW, Cm(0.6),
         "Official source, transparent methodology,"
         " honestly shown limitations.",
         sz=11, italic=True, color=ACCENT)

    _hline(sl, ML, Cm(2.0), CW)

    _img(sl, "21_quarterly_corridor.png", ML, Cm(2.3), COL_W)

    rx = ML + COL_W + Cm(1.0)

    _txt(sl, rx, Cm(2.5), COL_W, Cm(0.5),
         "What this brief shows", sz=10, bold=True, color=INK)

    _bullets(sl, rx, Cm(3.2), COL_W, Cm(4.0), [
        "\u25b8 Apartments in CHF/m\u00b2 (PPE),"
        " houses in CHF per dwelling"
        " \u2014 different units, never mixed.",

        "\u25b8 Median and quartiles (P25\u2013P75),"
        " not averages. Averages shown only"
        " where OCSTAT publishes no median.",

        "\u25b8 Communes with n<10 transactions"
        " are shown as grey, not zero."
        " Small samples are flagged.",
    ], sz=9)

    _txt(sl, rx, Cm(7.5), COL_W, Cm(0.5),
         "18 years of resilience", sz=10, bold=True, color=INK)

    _bullets(sl, rx, Cm(8.2), COL_W, Cm(5.0), [
        "\u25b8 Three macro shocks (2008, COVID,"
        " 2022 rate hikes) \u2014 no sustained"
        " median decline. Volumes absorb shocks,"
        " prices hold.",

        "\u25b8 House P75 surged +64% in 2021\u20132023"
        " then retreated \u2014 premium segment"
        " reacts, median absorbs.",

        "\u25b8 Apartment median continues to climb;"
        " house median plateaus in 2026."
        " Two markets are diverging.",
    ], sz=9)

    _txt(sl, rx, Cm(12.5), COL_W, Cm(0.5),
         "Source & reproducibility", sz=10, bold=True, color=INK)

    _bullets(sl, rx, Cm(13.2), COL_W, Cm(4.0), [
        "\u25b8 OCSTAT T 05.05 series \u2014 completed"
        " notarial transactions, Canton of Geneva."
        " Published by the cantonal statistics office.",

        "\u25b8 All charts generated from raw data"
        " via Python (pandas, matplotlib)."
        " Code and methodology available on GitHub.",
    ], sz=9)

    _hline(sl, ML, Cm(19.5), CW)
    _footer(sl)


def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page_price(prs)
    page_geography(prs)
    page_method(prs)

    out = OUT / "Geneva_Market_Brief_LinkedIn.pptx"
    prs.save(str(out))
    print(f"Saved -> {out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
